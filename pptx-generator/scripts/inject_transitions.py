# -*- coding: utf-8 -*-
"""python-pptx 原生注入页面切换动画（p:transition）+ round-trip 验证。
已实测可靠：不要用 zipfile 手动重写+字符串注入（PowerPoint 会报打不开）。
用法: python inject_transitions.py <src.pptx> <dst.pptx>
速度: slow/med/fast；效果: zoom-in|dissolve|push-l|push-r|wipe-l|wipe-r|split-horz|comb-l|wheel|fade
"""
import sys
from pptx import Presentation
from lxml import etree

P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

TRANSITIONS = [  # 默认 10 页示例：封面/结尾 slow，内容页 med，效果各不相同
    ('slow', 'zoom-in'), ('med', 'dissolve'), ('med', 'push-l'), ('med', 'wipe-l'),
    ('med', 'split-horz'), ('slow', 'comb-l'), ('med', 'wheel'), ('med', 'wipe-r'),
    ('med', 'push-l'), ('slow', 'fade'),
]

BODIES = {
    'zoom-in': '<p:zoom dir="in"/>', 'dissolve': '<p:dissolve/>',
    'push-l': '<p:push dir="l"/>', 'push-r': '<p:push dir="r"/>',
    'wipe-l': '<p:wipe dir="l"/>', 'wipe-r': '<p:wipe dir="r"/>',
    'split-horz': '<p:split orient="horz" dir="out"/>',
    'comb-l': '<p:comb dir="l"/>', 'wheel': '<p:wheel/>', 'fade': '<p:fade/>',
}

def build_tr(speed, kind):
    xml = f'<p:transition xmlns:p="{P}" spd="{speed}" advClick="1">{BODIES[kind]}</p:transition>'
    return etree.fromstring(xml.encode('utf-8'))

src = sys.argv[1] if len(sys.argv) > 1 else r'D:\Temp\deck.pptx'
dst = sys.argv[2] if len(sys.argv) > 2 else r'D:\Temp\deck_transitions.pptx'

prs = Presentation(src)
slides = list(prs.slides)
for i, slide in enumerate(slides):
    if i >= len(TRANSITIONS):
        speed, kind = 'med', 'fade'
    else:
        speed, kind = TRANSITIONS[i]
    slide._element.append(build_tr(speed, kind))
    print(f'slide{i+1}: {kind} ({speed})')
prs.save(dst)
print('SAVED:', dst)

# round-trip 验证
prs2 = Presentation(dst)
ok = all(f'{{{P}}}transition' in [c.tag for c in s._element] for s in prs2.slides)
print('ROUND-TRIP CHECK:', 'ALL OK' if ok else 'PROBLEM')
sys.exit(0 if ok else 1)
